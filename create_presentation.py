from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PLUM = RGBColor(0x7B, 0x2D, 0x8E)
PLUM_LIGHT = RGBColor(0xA8, 0x55, 0xF7)
DARK_BG = RGBColor(0x0A, 0x0A, 0x0F)
CARD_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xCC)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
RED = RGBColor(0xEF, 0x44, 0x44)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_fill(slide, left, top, width, height, color, radius=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=PLUM_LIGHT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(1), color=PLUM_LIGHT):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_divider(slide, top):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), top, Inches(11.333), Pt(1))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x44)
    shp.line.fill.background()

def add_card(slide, left, top, width, height, title, body_lines, icon_text=None, accent=PLUM):
    card = add_shape_fill(slide, left, top, width, height, CARD_BG)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.45),
                 title, font_size=16, color=WHITE, bold=True)
    y = top + Inches(0.6)
    for line in body_lines:
        add_text_box(slide, left + Inches(0.3), y, width - Inches(0.6), Inches(0.35),
                     line, font_size=12, color=LIGHT_GRAY)
        y += Inches(0.28)


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
# Accent stripe at top
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(6)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = PLUM
slide.shapes[-1].line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
             'PLUMFOLIO', font_size=56, color=PLUM_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.0), Inches(10), Inches(0.8),
             'A Smart Personal Finance Tracker', font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
add_divider(slide, Inches(3.9))
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
             'Final Year Project Presentation', font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
             'Leon Maunge', font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             'Supervisor: Dr. Shree Om  |  University of Botswana  |  2026',
             font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: AGENDA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Presentation Outline', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

agenda = [
    ('01', 'Goal & Mission', 'What Plumfolio aims to achieve'),
    ('02', 'Problem Definition', 'The gap we are addressing'),
    ('03', 'Solution & Features', 'How Plumfolio solves the problem'),
    ('04', 'System Architecture', 'Tech stack & design decisions'),
    ('05', 'User Interface', 'Intuitive design walkthrough'),
    ('06', 'Core Functionalities', 'Live demonstration of key features'),
    ('07', 'Validation & Security', 'Data integrity & protection measures'),
    ('08', 'Conclusion & Future Work', 'Summary and next steps'),
]

for i, (num, title, desc) in enumerate(agenda):
    row = i // 2
    col = i % 2
    x = Inches(1.2 + col * 5.6)
    y = Inches(2.0 + row * 1.2)
    add_shape_fill(slide, x, y, Inches(5), Inches(0.95), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.6), Inches(0.4),
                 num, font_size=20, color=PLUM_LIGHT, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.1), Inches(3.8), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.5), Inches(3.8), Inches(0.35),
                 desc, font_size=12, color=LIGHT_GRAY)


# ==================== SLIDE 3: GOAL / MISSION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Goal & Mission', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

add_text_box(slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(0.8),
             '"To empower individuals to take control of their personal finances through an '
             'intuitive, intelligent, and accessible web application."',
             font_size=20, color=PLUM_LIGHT, bold=False)

goals = [
    'Provide a centralised platform to track income, expenses, budgets, and savings goals',
    'Deliver real-time financial insights and predictive analytics to guide better decisions',
    'Enable receipt scanning and bank statement imports to reduce manual data entry',
    'Support multi-currency display for users across different regions',
    'Ensure data security through Supabase authentication and row-level security',
    'Create an intuitive, responsive interface that works on desktop and mobile',
]

add_bullet_list(slide, Inches(1.2), Inches(2.8), Inches(10.5), Inches(4), goals, font_size=17)


# ==================== SLIDE 4: PROBLEM DEFINITION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Problem Definition', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

problems = [
    ('Lack of Financial Awareness',
     'Many individuals lack visibility into where their money goes each month, '
     'leading to overspending, missed savings targets, and financial stress.'),
    ('Fragmented Tools',
     'Existing solutions are either too complex (enterprise software), too limited '
     '(basic spreadsheets), or locked behind expensive subscriptions.'),
    ('Manual Data Entry Burden',
     'Manually recording every transaction is tedious and error-prone, causing users '
     'to abandon financial tracking altogether.'),
    ('No Actionable Insights',
     'Most trackers only show numbers — they do not explain spending patterns, '
     'detect anomalies, or forecast future cash flow.'),
]

for i, (title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = Inches(1.0 + col * 5.8)
    y = Inches(1.8 + row * 2.5)
    add_card(slide, x, y, Inches(5.4), Inches(2.1), title, [desc],
             accent=[RED, AMBER, PLUM_LIGHT, RGBColor(0x38, 0xBD, 0xF8)][i])


# ==================== SLIDE 5: SOLUTION & FEATURES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Solution & Key Features', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

features = [
    ('Dashboard', 'Real-time financial overview with\nmonthly summaries & charts', GREEN),
    ('Transactions', 'Full CRUD + CSV import +\nreceipt scanning (OCR)', PLUM_LIGHT),
    ('Budgets & Goals', 'Category budgets, formula allocator,\nsavings goals with progress', AMBER),
    ('Analytics', 'Trend analysis with interactive\nChart.js visualisations', RGBColor(0x38, 0xBD, 0xF8)),
    ('AI Insights', 'Rule-based engine: spending patterns,\nanomalies, forecasts, tips', RGBColor(0xEC, 0x48, 0x99)),
    ('Reports', 'Period comparisons, PDF export,\nyear-over-year analysis', RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, accent) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 2.7)
    add_card(slide, x, y, Inches(3.6), Inches(2.2), title, desc.split('\n'), accent=accent)


# ==================== SLIDE 6: SYSTEM ARCHITECTURE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'System Architecture & Tech Stack', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

# Left: stack
stack_items = [
    ('Frontend', 'React 18.3 + React Router v6\nChart.js + Lucide Icons\nResponsive CSS (mobile-first)', PLUM_LIGHT),
    ('Backend', 'Supabase (PostgreSQL)\nREST API + Realtime subscriptions\nEdge Functions (receipt scanning)', GREEN),
    ('Authentication', 'Supabase Auth (JWT)\nEmail verification + password recovery\nRow-Level Security (RLS)', AMBER),
    ('Deployment', 'GitHub Pages (SPA)\nPWA manifest for installability\nCDN-served static assets', RGBColor(0x38, 0xBD, 0xF8)),
]

for i, (title, desc, accent) in enumerate(stack_items):
    y = Inches(1.8 + i * 1.3)
    add_shape_fill(slide, Inches(1), y, Inches(5.5), Inches(1.1), CARD_BG)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y, Pt(5), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(2), Inches(0.35),
                 title, font_size=15, color=accent, bold=True)
    for j, line in enumerate(desc.split('\n')):
        add_text_box(slide, Inches(1.3), y + Inches(0.4 + j * 0.23), Inches(4.8), Inches(0.25),
                     line, font_size=11, color=LIGHT_GRAY)

# Right: architecture flow
add_text_box(slide, Inches(7.2), Inches(1.8), Inches(5), Inches(0.4),
             'Architecture Flow', font_size=18, color=WHITE, bold=True)

flow_items = [
    ('User', 'Browser / Mobile', PLUM_LIGHT),
    ('React SPA', 'Components + Context API', WHITE),
    ('Supabase Client', 'Auth + REST + Realtime', GREEN),
    ('PostgreSQL', 'Tables + RLS Policies', AMBER),
    ('Edge Functions', 'Receipt OCR Processing', RGBColor(0x38, 0xBD, 0xF8)),
]

for i, (title, sub, color) in enumerate(flow_items):
    y = Inches(2.4 + i * 0.95)
    add_shape_fill(slide, Inches(7.5), y, Inches(4.5), Inches(0.7), CARD_BG)
    add_text_box(slide, Inches(7.8), y + Inches(0.05), Inches(2), Inches(0.3),
                 title, font_size=14, color=color, bold=True)
    add_text_box(slide, Inches(9.8), y + Inches(0.05), Inches(2), Inches(0.3),
                 sub, font_size=11, color=LIGHT_GRAY)
    if i < len(flow_items) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.5), y + Inches(0.7),
                                        Inches(0.35), Inches(0.22))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x66)
        arrow.line.fill.background()


# ==================== SLIDE 7: USER INTERFACE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'User Interface Design', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

ui_principles = [
    ('Dark Theme', 'Reduces eye strain for daily financial\nreview; modern aesthetic with plum accents'),
    ('Responsive Layout', 'Desktop sidebar + mobile bottom navigation;\nworks on all screen sizes'),
    ('Intuitive Navigation', '6 main sections accessible from sidebar;\nclear visual hierarchy per page'),
    ('Real-time Feedback', 'Toast notifications, inline validation,\nloading states for every action'),
    ('Accessibility', 'High contrast text, semantic HTML,\nkeyboard-navigable forms'),
    ('Progressive Web App', 'Installable on mobile devices;\nnative app-like experience'),
]

for i, (title, desc) in enumerate(ui_principles):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 2.7)
    add_card(slide, x, y, Inches(3.6), Inches(2.2), title, desc.split('\n'), accent=PLUM_LIGHT)


# ==================== SLIDE 8: CORE FUNCTIONALITIES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Core Functionalities (Demo)', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
             'The following features will be demonstrated live:', font_size=16, color=LIGHT_GRAY)

demo_items = [
    ('1. Authentication Flow', 'Sign up with validation \u2192 email verification \u2192 sign in \u2192 session management'),
    ('2. Transaction Management', 'Add income/expense \u2192 edit \u2192 delete \u2192 filter by month/year \u2192 search'),
    ('3. Receipt Scanning', 'Upload receipt image \u2192 OCR processing \u2192 auto-populated form fields'),
    ('4. Bank Statement Import', 'Upload CSV \u2192 auto-categorisation \u2192 duplicate detection \u2192 bulk import'),
    ('5. Budget Management', 'Create budgets per category \u2192 formula allocator \u2192 over-budget warnings'),
    ('6. Savings Goals', 'Set targets \u2192 track progress \u2192 contribute funds \u2192 goal completion'),
    ('7. Analytics & Reports', 'Interactive charts \u2192 period comparisons \u2192 PDF export'),
    ('8. AI Insights', 'Spending pattern analysis \u2192 anomaly detection \u2192 forecasting \u2192 savings tips'),
]

for i, (title, desc) in enumerate(demo_items):
    y = Inches(2.1 + i * 0.62)
    add_shape_fill(slide, Inches(1), y, Inches(11.333), Inches(0.52), CARD_BG)
    add_text_box(slide, Inches(1.3), y + Inches(0.06), Inches(3), Inches(0.4),
                 title, font_size=13, color=PLUM_LIGHT, bold=True)
    add_text_box(slide, Inches(4.3), y + Inches(0.06), Inches(7.5), Inches(0.4),
                 desc, font_size=12, color=LIGHT_GRAY)


# ==================== SLIDE 9: VALIDATION & SECURITY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Validation & Security', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

sec_left = [
    ('Input Validation', [
        'All amounts: positive, max 10M, 2 decimal places',
        'Descriptions: 2\u2013100 chars, XSS pattern rejection',
        'Passwords: min 8 chars, uppercase + lowercase + number',
        'Emails: format check, length limit (254), trimmed',
        'Dates: no future dates, max 5 years back',
        'Budgets: min P10, month range validation',
    ]),
    ('Duplicate Prevention', [
        'Transactions: fuzzy matching on amount + description + date',
        'Budgets: one per category per month enforced',
        'Sign-up: existing email detection (including Supabase silent duplicates)',
    ]),
]

sec_right = [
    ('Authentication & Authorisation', [
        'Supabase Auth with JWT session tokens',
        'Email verification before first sign-in',
        'Password recovery via secure reset links',
        'Row-Level Security (RLS) on all database tables',
        'Protected routes redirect unauthenticated users',
    ]),
    ('Data Protection', [
        'HTTPS encryption in transit',
        'Public anon key + server-side RLS enforcement',
        'Account deletion: purges all user data',
        'No sensitive data stored in localStorage',
        'Control character & script injection blocking',
    ]),
]

for col_idx, items in enumerate([sec_left, sec_right]):
    x_base = Inches(0.8 + col_idx * 6.0)
    y_pos = Inches(1.7)
    for title, bullets in items:
        add_text_box(slide, x_base + Inches(0.2), y_pos, Inches(5.4), Inches(0.4),
                     title, font_size=17, color=PLUM_LIGHT, bold=True)
        y_pos += Inches(0.45)
        for bullet in bullets:
            add_text_box(slide, x_base + Inches(0.4), y_pos, Inches(5.2), Inches(0.3),
                         f'\u2022  {bullet}', font_size=12, color=LIGHT_GRAY)
            y_pos += Inches(0.28)
        y_pos += Inches(0.2)


# ==================== SLIDE 10: FUNCTIONAL REQUIREMENTS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Functional Requirements Coverage', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

frs = [
    ('FR-1', 'User Authentication', 'Sign up, sign in, email verification, password recovery', GREEN),
    ('FR-2', 'Transaction Management', 'CRUD operations, filtering, search, categorisation', GREEN),
    ('FR-3', 'Budget Management', 'Create, edit, delete budgets; formula-based allocation', GREEN),
    ('FR-4', 'Savings Goals', 'Set targets, track progress, fund contributions', GREEN),
    ('FR-5', 'Financial Insights', 'Rule-based analysis: patterns, anomalies, tips', GREEN),
    ('FR-6', 'Predictive Analytics', 'Weighted moving averages, trend detection, forecasts', GREEN),
    ('FR-7', 'Receipt Scanning', 'OCR with Tesseract.js + Edge Function fallback', GREEN),
    ('FR-8', 'Data Import', 'CSV parsing, auto-categorisation, duplicate detection', GREEN),
    ('FR-9', 'Reports & Export', 'Period reports, comparisons, PDF export', GREEN),
    ('FR-10', 'Multi-Currency', 'Live FX rates, user-selectable display currency', GREEN),
]

for i, (code, name, desc, status) in enumerate(frs):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(1.7 + row * 1.05)
    add_shape_fill(slide, x, y, Inches(5.6), Inches(0.85), CARD_BG)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.65), Inches(0.3),
                 code, font_size=11, color=PLUM_LIGHT, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.08), Inches(3.5), Inches(0.3),
                 name, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.4), Inches(4.2), Inches(0.35),
                 desc, font_size=10, color=LIGHT_GRAY)
    add_text_box(slide, x + Inches(4.8), y + Inches(0.08), Inches(0.6), Inches(0.3),
                 '\u2713', font_size=16, color=status, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 11: CHALLENGES & LESSONS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Challenges & Lessons Learned', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

challenges = [
    ('Supabase RLS & Auth Edge Cases',
     'Handling silent duplicate emails, session persistence across tabs, and '
     'designing RLS policies that balance security with usability.'),
    ('Receipt OCR Accuracy',
     'Raw OCR text is noisy — built two parsers (heuristic + strict) and '
     'image preprocessing (adaptive threshold, contrast enhancement) to improve accuracy.'),
    ('Budget Logic Complexity',
     'Balancing user expectations (income vs. balance-based budgeting) with correct '
     'financial calculations required iterative UX testing and logic refinement.'),
    ('Client-Side Performance',
     'Running insights, predictions, and cross-reference engines in the browser '
     'required careful optimisation and 30-second polling with Realtime fallback.'),
]

for i, (title, desc) in enumerate(challenges):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 5.9)
    y = Inches(1.7 + row * 2.6)
    add_card(slide, x, y, Inches(5.5), Inches(2.2), title, [desc], accent=AMBER)


# ==================== SLIDE 12: FUTURE WORK ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Future Work & Enhancements', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

future = [
    ('Direct Bank Integration', 'Connect to banking APIs for automatic\ntransaction syncing (e.g., Plaid, Stitch)'),
    ('Machine Learning Insights', 'Replace rule-based engine with ML\nmodels for personalised predictions'),
    ('Mobile Native App', 'React Native version with push\nnotifications and biometric auth'),
    ('Collaborative Budgets', 'Shared household budgets with\nmulti-user access and permissions'),
    ('Recurring Transactions', 'Auto-detect and auto-create recurring\nincome and expense entries'),
    ('Investment Tracking', 'Portfolio tracking, stock prices,\nnet worth dashboard'),
]

for i, (title, desc) in enumerate(future):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.7 + row * 2.7)
    add_card(slide, x, y, Inches(3.6), Inches(2.2), title, desc.split('\n'),
             accent=RGBColor(0x38, 0xBD, 0xF8))


# ==================== SLIDE 13: CONCLUSION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
             'Conclusion', font_size=36, color=WHITE, bold=True)
add_accent_bar(slide, Inches(1), Inches(1.25), Inches(2), Pt(4), PLUM_LIGHT)

conclusions = [
    'Plumfolio successfully addresses the problem of personal finance management by providing '
    'a comprehensive, intuitive, and intelligent web application.',
    'All 10 functional requirements have been implemented and are fully operational.',
    'The system combines modern web technologies (React 18, Supabase, Chart.js) with '
    'rule-based AI engines for insights and predictions.',
    'Receipt scanning and bank statement imports significantly reduce the manual data entry burden.',
    'Comprehensive validation and security measures ensure data integrity and user protection.',
    'The responsive design and PWA capabilities make Plumfolio accessible across all devices.',
]

for i, item in enumerate(conclusions):
    y = Inches(1.7 + i * 0.85)
    add_shape_fill(slide, Inches(1), y, Inches(11.333), Inches(0.7), CARD_BG)
    add_text_box(slide, Inches(1.3), y + Inches(0.12), Inches(10.7), Inches(0.5),
                 f'\u2713  {item}', font_size=14, color=LIGHT_GRAY)


# ==================== SLIDE 14: THANK YOU / Q&A ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Pt(6)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = PLUM
slide.shapes[-1].line.fill.background()

add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1),
             'Thank You', font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             'Questions & Discussion', font_size=28, color=PLUM_LIGHT, alignment=PP_ALIGN.CENTER)
add_divider(slide, Inches(4.2))
add_text_box(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.5),
             'Leon Maunge  |  University of Botswana  |  2026',
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.5),
             'leonfinalyearproject.github.io/Plumfolio',
             font_size=14, color=PLUM_LIGHT, alignment=PP_ALIGN.CENTER)


# ==================== SAVE ====================
output_path = r'c:\Users\Leonm\Music\Plumfolio-main\Plumfolio_Presentation.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
